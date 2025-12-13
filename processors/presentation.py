"""Presentation (PowerPoint) processing utilities."""

import os
from pptx import Presentation
from pptx.util import Inches
import io


def get_slide_text(slide):
    """Extract all text from a slide."""
    texts = []

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())

        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells]
                texts.append(' | '.join(row_texts))

    return '\n'.join(texts)


def get_slide_notes(slide):
    """Extract speaker notes from slide."""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            return notes_slide.notes_text_frame.text.strip()
    return ""


def get_slide_title(slide):
    """Extract slide title."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    return ""


def get_slide_image_ratio(slide):
    """Calculate ratio of image shapes to total shapes."""
    total_shapes = len(slide.shapes)
    if total_shapes == 0:
        return 0

    image_shapes = sum(1 for shape in slide.shapes if shape.shape_type == 13)
    return image_shapes / total_shapes


def should_use_ocr(prs, use_ocr=None):
    """
    Determine if OCR should be used.

    Args:
        prs: Presentation object
        use_ocr: None (auto), True (force), False (skip)

    Returns:
        (should_ocr, reason)
    """
    if use_ocr is True:
        return True, "User requested OCR"

    if use_ocr is False:
        return False, "OCR disabled by user"

    total_chars = 0
    total_slides = len(prs.slides)
    high_image_slides = 0

    for slide in prs.slides:
        text = get_slide_text(slide)
        total_chars += len(text)

        img_ratio = get_slide_image_ratio(slide)
        if img_ratio > 0.8:
            high_image_slides += 1

    chars_per_slide = total_chars / total_slides if total_slides > 0 else 0

    if chars_per_slide < 50:
        return True, f"Low text density ({chars_per_slide:.0f} chars/slide)"

    if high_image_slides / total_slides > 0.8:
        return True, f"Image-heavy presentation ({high_image_slides}/{total_slides} slides)"

    return False, "Sufficient text content"


def extract_text_direct(prs):
    """Extract text directly from presentation."""
    slides_data = []
    md_parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = get_slide_title(slide)
        content = get_slide_text(slide)
        notes = get_slide_notes(slide)

        slides_data.append({
            "slide_number": slide_num,
            "title": title,
            "text": content,
            "notes": notes,
            "char_count": len(content)
        })

        slide_md = f"## Slide {slide_num}"
        if title:
            slide_md += f": {title}"
        slide_md += "\n\n"

        if content:
            slide_md += content + "\n"

        if notes:
            slide_md += f"\n> **Speaker Notes:** {notes}\n"

        md_parts.append(slide_md)

    full_text = "# Presentation\n\n---\n\n" + "\n---\n\n".join(md_parts)
    return full_text, slides_data


def extract_text_with_ocr(prs, file_path):
    """Extract text using OCR for each slide."""
    from processors import gemini_ocr
    import pymupdf

    slides_data = []
    md_parts = []

    pdf_path = file_path.replace('.pptx', '_temp.pdf').replace('.ppt', '_temp.pdf')
    try:
        from subprocess import run, PIPE
        run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir',
             os.path.dirname(file_path), file_path], stdout=PIPE, stderr=PIPE)
    except:
        pass

    if os.path.exists(pdf_path):
        doc = pymupdf.open(pdf_path)
        for page_num, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2))
            img_bytes = pix.tobytes("png")

            try:
                slide_text = gemini_ocr.extract_text_from_image(img_bytes, 'image/png')
            except Exception as e:
                slide_text = f"[OCR Error: {str(e)}]"

            slides_data.append({
                "slide_number": page_num,
                "title": "",
                "text": slide_text,
                "notes": "",
                "char_count": len(slide_text)
            })

            md_parts.append(f"## Slide {page_num}\n\n{slide_text}")

        doc.close()
        os.remove(pdf_path)
    else:
        return extract_text_direct(prs)

    full_text = "# Presentation\n\n---\n\n" + "\n---\n\n".join(md_parts)
    return full_text, slides_data


def process(file_path, use_ocr=None):
    """
    Extract text from presentation and convert to Markdown.

    Args:
        file_path: Path to .pptx file
        use_ocr: None (auto-detect), True (force OCR), False (direct only)

    Returns:
        {
            "text": str,
            "file_size": int,
            "num_pages": int,
            "processing_method": str,
            "metadata": dict
        }
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Presentation file not found: {file_path}")

    file_size = os.path.getsize(file_path)

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Invalid presentation file: {e}")

    num_slides = len(prs.slides)
    should_ocr, ocr_reason = should_use_ocr(prs, use_ocr)

    if should_ocr:
        full_text, slides_data = extract_text_with_ocr(prs, file_path)
        processing_method = 'ocr'
    else:
        full_text, slides_data = extract_text_direct(prs)
        processing_method = 'direct'

    if not full_text.strip():
        raise ValueError("Could not extract any text from presentation.")

    return {
        "text": full_text,
        "file_size": file_size,
        "num_pages": num_slides,
        "processing_method": processing_method,
        "metadata": {
            "ocr_reason": ocr_reason,
            "slides": slides_data,
            "extracted_slides": len([s for s in slides_data if s["char_count"] > 0])
        }
    }


def analyze_presentation(file_path):
    """
    Analyze presentation to determine if OCR is recommended.

    Args:
        file_path: Path to presentation

    Returns:
        {
            "num_pages": int,
            "ocr_recommended": bool,
            "ocr_reason": str
        }
    """
    try:
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        should_ocr, reason = should_use_ocr(prs, None)

        return {
            "num_pages": num_slides,
            "ocr_recommended": should_ocr,
            "ocr_reason": reason
        }
    except Exception as e:
        return {
            "num_pages": 0,
            "ocr_recommended": True,
            "ocr_reason": f"Analysis failed: {str(e)}"
        }
